import yaml
from google_images_search import GoogleImagesSearch
from pptx.shapes.placeholder import PicturePlaceholder


def get_images(presentation, name, safety=1):
    number_of_images = calculate_number_of_images_needed(presentation)

    dev_api_key, project_cx = get_gis_credentials()

    gis = GoogleImagesSearch(dev_api_key, project_cx)
    _search_params = {
        'q': name,
        'num': number_of_images + safety}

    gis.search(search_params=_search_params,
               path_to_dir='static/core/img', custom_image_name=name)


def calculate_number_of_images_needed(presentation):
    number_of_images_needed = 0

    for slide in presentation.slides:
        for placeholder in slide.placeholders:
            if isinstance(placeholder, PicturePlaceholder):
                number_of_images_needed += 1

    return number_of_images_needed


def get_gis_credentials():
    with open('static/core/extras/config.yaml', 'r') as stream:
        try:
            config = yaml.safe_load(stream)
        except yaml.YAMLError as exc:
            print(exc)

    return config['google_images_search_dev_api_key'], config['google_images_search_project_cx']
