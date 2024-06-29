import torch
from os import remove as os_remove
from glob import glob
from pickle import load as pickle_load
from random import randint
from transformers import BertTokenizer
from pptx import Presentation
from pptx.shapes.placeholder import PicturePlaceholder
from core.language_service import get_language, translator, translate_text
from core.image_service import get_images
from core.summary_service import summarize_text
from core.extract_content_service import create_content_tree
from core.wikipedia_service import find_url

themes = {0: "Language", 1: "Geography", 2: "History", 3: "Science", 4: "Biology", 5: "Sport", 6: "Technology",
          7: "Business", 8: "Space", 9: "Art", 10: "Other"}

theme_model = pickle_load(
    open('static/core/extras/theme_classifier_ml_base.pickle', 'rb'))

bert_model = 'bert-base-uncased'
app_name = 'new_presentations'

tokenizer = BertTokenizer.from_pretrained(bert_model)


def main(requested_title, requested_language, requested_number_of_slides, user_full_name, requested_min_words=100, requested_max_words=120):
    url, summary, sections, final_title = find_url(
        requested_title, requested_language)

    text = create_content_tree(summary, sections, final_title)

    summarized_text = summarize_text(
        text, requested_number_of_slides, final_title, requested_min_words, requested_max_words)

    final_text = translate_text(summarized_text, requested_language)

    presentation = make_presentation(
        final_text, requested_title, requested_language, user_full_name)

    return presentation


def make_presentation(final_text, title, language, user_full_name):

    presentation, category = create_presentation_from_template(
        final_text, title)

    main_titles, sub_titles, paragraphs, new_title = unpack_content(
        final_text, language, title, category, user_full_name)

    populate_text(presentation, main_titles, sub_titles, paragraphs)

    remove_unused_slides(presentation, len(paragraphs))

    populate_images(presentation, title)

    populate_thank_you_slide(presentation, get_language(language))

    new_title = format_title(new_title)
    final_name = f'static/core/presentations/{app_name}_{new_title}.pptx'
    presentation.save(final_name)

    return f'{app_name}_{new_title}.pptx'


def create_presentation_from_template(final_text, title):
    category = get_template_category(
        final_text['About the ' + title]['major subject'])

    template = category + str(randint(1, 3))

    presentation = Presentation(
        'static/core/presentation_templates/' + template + '.pptx')

    return presentation, category


def unpack_content(final_text, language, title, category, user_full_name):
    lang = get_language(language)

    if lang != 'en':
        new_title = translator.translate(title, src='en', dest=lang).text
        new_category = translator.translate(category, src='en', dest=lang).text
    else:
        new_title = title
        new_category = category

    main_titles = [new_title]
    sub_titles = [new_category]
    paragraphs = [user_full_name]

    for key in final_text:
        main_titles.append(key)
        for sub_key in final_text[key]:
            sub_titles.append(sub_key)
            paragraphs.append(final_text[key][sub_key])
        for _ in range(len(sub_titles) - len(main_titles)):
            main_titles.append('')

    return main_titles, sub_titles, paragraphs, new_title


def populate_text(presentation, main_titles, sub_titles, paragraphs):
    for i in range(min(len(paragraphs), 20)):
        slide = presentation.slides[i]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            if text_frame.text.startswith('t1'):
                active_main_title = main_titles[i] if main_titles[i] else active_main_title
                run = text_frame.paragraphs[0].runs[0]
                run.text = active_main_title
            elif text_frame.text.startswith('t2'):
                active_sub_title = sub_titles[i] if sub_titles[i] else active_sub_title
                run = text_frame.paragraphs[0].runs[0]
                run.text = active_sub_title
            elif text_frame.text.startswith('#t3#'):
                active_paragraph = paragraphs[i] if paragraphs[i] else active_paragraph
                run = text_frame.paragraphs[0].runs[0]
                run.text = active_paragraph


def remove_unused_slides(presentation, num_paragraphs):
    # Needs to be done like this because presentation.sledes doesn't have delete method in version 0.6.21
    while len(presentation.slides) - 1 > num_paragraphs:
        xml_slides = presentation.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[-2])


def populate_images(presentation, title):
    get_images(presentation, title)

    insert_images(presentation)

    remove_used_image_files(glob('static/core/img/*'))


def populate_thank_you_slide(presentation, language):
    thank_you = "Thank you for your attention"

    if language != 'en':
        thank_you = translator.translate(
            thank_you, src='en', dest=language).text

    last_slide = presentation.slides[-1]

    for shape in last_slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        if text_frame.text.startswith('t*'):
            run = text_frame.paragraphs[0].runs[0]
            run.text = thank_you


def get_template_category(text):
    test_dataset = [text]

    test_loader = torch.utils.data.DataLoader(
        test_dataset, batch_size=1, collate_fn=pad)

    return themes[test_loop(theme_model, test_loader)]


def insert_images(presentation):
    files = glob('static/core/img/*')
    safety = True
    image_serial_number = 0

    for slide in presentation.slides:
        for placeholder in slide.placeholders:
            if isinstance(placeholder, PicturePlaceholder):
                try:
                    placeholder.insert_picture(files[image_serial_number])
                except Exception:
                    if safety:
                        safety = False
                        image_serial_number += 1
                        try:
                            placeholder.insert_picture(
                                files[image_serial_number])
                        except Exception:
                            pass
                image_serial_number += 1


def test_loop(model, dataloader):
    model.eval()
    for occurrence in dataloader:
        out = model(occurrence)
        labs = out.logits.argmax(dim=1)
        return labs[0].item()


def pad(b):
    v = [tokenizer.encode(x) for x in b]
    l = max(map(len, v))

    return torch.stack(
        [torch.nn.functional.pad(torch.tensor(t), (0, l - len(t)), mode='constant', value=0) for t in v])


def remove_used_image_files(files):
    for file in files:
        os_remove(file)


def format_title(title):
    return title.replace(' ', '_')
